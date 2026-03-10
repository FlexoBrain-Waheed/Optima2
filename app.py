import streamlit as st
from pptx import Presentation
import io

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Strategic Investment Evaluation"
    slide1.placeholders[1].text = "SOMA Optima 2 vs. W&H Alphaflex\n\nPresented by Flexo Consultation Services | Waheed Alkarraein"

    # Slide 2: Market Positioning
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "The Paradigm Shift: Flagship vs. Economy"
    slide2.placeholders[1].text = (
        "• W&H Alphaflex (Economy Class): Positioned as an 'Entry-Level' press. Stripped down to fit lower budgets.\n\n"
        "• SOMA Optima 2 (Flagship Class): The absolute pinnacle of SOMA's technology. A fully-equipped, premium machine.\n\n"
        "• The Strategic Question: Does it make sense to pay over €2 Million for a stripped-down economy model when you can own top-tier technology for the same budget?"
    )

    # Slide 3: Technological Maturity
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Proven Reliability vs. The Testing Ground"
    slide3.placeholders[1].text = (
        "• SOMA Optima 2: A mature, 100% reliable platform. Built on deep operational experience with zero structural flaws.\n\n"
        "• W&H Alphaflex: Newly launched in mid-2024. Its economical structure is largely untested in harsh, long-term operational environments.\n\n"
        "• Conclusion: Never let your factory be the testing ground for a new economy model."
    )

    # Slide 4: Core Technology & Performance
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Smart Engineering for Maximum Profit"
    slide4.placeholders[1].text = (
        "• SOMA Optima 2 (Bounce Killer): Uses aerospace-grade Carbon Composite mandrels and integrated Ink Thermal Stabilization (ITS).\n\n"
        "• W&H Alphaflex: Relies on heavy steel mandrels and lacks an integrated cooling system, leading to potential bouncing and color instability.\n\n"
        "• Zero Waste Guarantee: SOMA's IRIS topography-based setup drastically reduces setup time and material waste compared to semi-automatic systems."
    )

    # Slide 5: Operator Ergonomics
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Empowering Your Workforce"
    slide5.placeholders[1].text = (
        "• SOMA Optima 2: Features 19 kg Lightweight Air Shafts.\n\n"
        "• W&H Alphaflex: Uses standard, heavy mechanical shafts.\n\n"
        "• The Result: SOMA dramatically reduces physical strain on operators, ensuring faster, safer, and much more efficient manual roll changes every single shift."
    )

    # Slide 6: Total Cost of Ownership
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "The Hidden Costs of 'Base' Offers"
    slide6.placeholders[1].text = (
        "• The SOMA Optima 2 is a complete 'Turn-Key' solution ready for immediate high-quality production.\n\n"
        "• The competitor's offer requires over €200,000 in hidden additions (Cooling, Cameras, Sleeves, Carbon Shafts) to reach the same operational level.\n\n"
        "• Financial Flexibility: SOMA offers flexible wire transfers (Net), whereas W&H requires a 75% Letter of Credit (L/C), tying up your cash flow and adding heavy bank fees."
    )

    # Slide 7: Partnership & Challenge
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "A Complete Partnership & The SOMA Challenge"
    slide7.placeholders[1].text = (
        "• Exclusive Turn-Key Package: Special massive discount if investing in the complete suite (Press, S-Mount, Lamiflex, Pluto).\n\n"
        "• Unbeatable Training: 10 days on-site + 2 extra weeks later + 1 Full Month in Europe at the SOMA factory for your team.\n\n"
        "• The SOMA Challenge: We invite you to our factory. Bring your most complex printing designs. Watch us set it up and print HD quality with near-zero waste in minutes. Trust what you see, not just ink on paper."
    )
    
    # Save the presentation to a BytesIO object
    binary_output = io.BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

# Streamlit App UI
st.set_page_config(page_title="Flexo Presentation Generator", page_icon="📊")

st.title("Flexo Consultation Services")
st.subheader("Presentation Generator")
st.write("Click the button below to generate and download the strategic comparison presentation.")

if st.button("Generate PowerPoint Presentation"):
    pptx_data = create_presentation()
    st.download_button(
        label="📥 Download .pptx File",
        data=pptx_data,
        file_name="Strategic_Investment_Evaluation_SOMA_vs_WH.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    st.success("Presentation generated successfully! Click the download button above.")
